from PIL import Image


def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    
    im = Image.open(image_url)
    width, height = im.size

    
    placeholder.height = height
    placeholder.width = width

    
    placeholder = placeholder.insert_picture(image_url)

    
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side
from pptx import Presentation
import os

prs = Presentation()

layout8 = prs.slide_layouts[8];
slide = prs.slides.add_slide(layout8);

title = slide.shapes.title.text = "Assignment";
sub = slide.placeholders[2].text = "watermark1 image";
_add_image(slide,1,"watermark1.jpg");

layout82 = prs.slide_layouts[8];
slide1 = prs.slides.add_slide(layout82);

title1 = slide.shapes.title.text = "Assignment";
sub1 = slide.placeholders[2].text = "watermark2 image";
_add_image(slide1,1,"watermark2.jpg");

layout83 = prs.slide_layouts[8];
slide2 = prs.slides.add_slide(layout83);

title2 = slide.shapes.title.text = "Assignment";
sub2 = slide.placeholders[2].text = "watermark3 image";
_add_image(slide2,1,"watermark3.jpg");

layout84 = prs.slide_layouts[8];
slide3 = prs.slides.add_slide(layout84);

title3 = slide.shapes.title.text = "Assignment";
sub3 = slide.placeholders[2].text = "watermark4 image";
_add_image(slide3,1,"watermark4.jpg");

layout85 = prs.slide_layouts[8];
slide4 = prs.slides.add_slide(layout85);

title4 = slide.shapes.title.text = "Assignment";
sub4 = slide.placeholders[2].text = "watermark5 image";
_add_image(slide4,1,"watermark5.jpg");

prs.save("ppt_assignment.pptx")
os.startfile("ppt_assignment.pptx")